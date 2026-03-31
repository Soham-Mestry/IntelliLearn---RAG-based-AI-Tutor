"""
RAG (Retrieval Augmented Generation) pipeline using LangChain.
Handles document processing (PDF/DOCX/PPT), embedding generation, and Q&A.
"""
import os
import pdfplumber
from docx import Document as DocxDocument
from pptx import Presentation
from langchain_text_splitters import RecursiveCharacterTextSplitter
from langchain_community.embeddings import HuggingFaceEmbeddings
from langchain_google_genai import ChatGoogleGenerativeAI
from langchain_core.prompts import PromptTemplate
from sqlalchemy.orm import Session
from sqlalchemy import text
from models import Embedding, Note
import numpy as np


# Initialize HuggingFace embeddings model (384 dimensions)
embeddings_model = HuggingFaceEmbeddings(
    model_name="sentence-transformers/all-MiniLM-L6-v2",
    model_kwargs={'device': 'cpu'},
    encode_kwargs={'normalize_embeddings': True}
)

# Initialize Google Gemini LLM
llm = ChatGoogleGenerativeAI(
    model="gemini-2.5-flash-lite",
    google_api_key=os.getenv("GOOGLE_API_KEY"),
    temperature=0.3
)

# Text splitter for chunking documents
text_splitter = RecursiveCharacterTextSplitter(
    chunk_size=500,
    chunk_overlap=50,
    length_function=len,
    separators=["\n\n", "\n", " ", ""]
)


def extract_text_from_pdf(filepath: str) -> str:
    """Extract text from PDF file using pdfplumber"""
    text = ""
    try:
        with pdfplumber.open(filepath) as pdf:
            for page in pdf.pages:
                page_text = page.extract_text()
                if page_text:
                    text += page_text + "\n"
    except Exception as e:
        raise ValueError(f"Error extracting PDF: {str(e)}")
    
    return text.strip()


def extract_text_from_docx(filepath: str) -> str:
    """Extract text from DOCX file using python-docx"""
    text = ""
    try:
        doc = DocxDocument(filepath)
        for paragraph in doc.paragraphs:
            text += paragraph.text + "\n"
    except Exception as e:
        raise ValueError(f"Error extracting DOCX: {str(e)}")
    
    return text.strip()


def extract_text_from_ppt(filepath: str) -> str:
    """Extract text from PPTX file using python-pptx (only supports PPTX, not old PPT format)"""
    text = ""
    try:
        prs = Presentation(filepath)
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    text += shape.text + "\n"
    except Exception as e:
        error_msg = str(e)
        # Check if it's an old PPT format issue
        if "relationship" in error_msg.lower() or "officeDocument" in error_msg:
            raise ValueError("Only PPTX format is supported. Please convert old PPT files to PPTX format and try again.")
        else:
            raise ValueError(f"Error extracting PowerPoint: {error_msg}")
    
    return text.strip()


def process_and_embed_document(filepath: str, note_id: str, db: Session) -> int:
    """
    Process document, split into chunks, generate embeddings, and store in database.
    Returns the number of chunks created.
    """
    # Determine file type and extract text
    file_extension = filepath.lower().split('.')[-1]
    
    if file_extension == 'pdf':
        text = extract_text_from_pdf(filepath)
    elif file_extension in ['docx', 'doc']:
        text = extract_text_from_docx(filepath)
    elif file_extension in ['ppt', 'pptx']:
        text = extract_text_from_ppt(filepath)
    else:
        raise ValueError(f"Unsupported file type: {file_extension}")
    
    if not text or len(text.strip()) == 0:
        raise ValueError("No text could be extracted from the document")
    
    # Split text into chunks
    chunks = text_splitter.split_text(text)
    
    if not chunks:
        raise ValueError("No chunks created from document")
    
    # Generate embeddings and store in database
    chunk_count = 0
    for chunk in chunks:
        if chunk.strip():  # Only process non-empty chunks
            # Generate embedding vector
            embedding_vector = embeddings_model.embed_query(chunk)
            
            # Create embedding record
            embedding_record = Embedding(
                note_id=note_id,
                text_chunk=chunk,
                embedding=embedding_vector
            )
            db.add(embedding_record)
            chunk_count += 1
    
    db.commit()
    return chunk_count


def similarity_search(query: str, db: Session, top_k: int = 3, subject_id: str = None) -> list:
    """
    Perform similarity search using pgvector.
    Optionally filters by subject_id to scope results to a specific subject.
    Returns top_k most similar text chunks with their context.
    """
    # Generate query embedding
    query_embedding = embeddings_model.embed_query(query)
    
    # Convert to PostgreSQL array format string
    embedding_str = "[" + ",".join(map(str, query_embedding)) + "]"
    
    # Build SQL query with optional subject filter
    if subject_id:
        sql_query = text("""
            SELECT 
                e.text_chunk,
                n.filename,
                s.name as subject_name,
                (1 - (e.embedding <=> CAST(:query_embedding AS vector))) as similarity
            FROM embeddings e
            JOIN notes n ON e.note_id = n.id
            JOIN subjects s ON n.subject_id = s.id
            WHERE s.id = CAST(:subject_id AS UUID)
            ORDER BY e.embedding <=> CAST(:query_embedding AS vector)
            LIMIT :top_k
        """)
        result = db.execute(
            sql_query,
            {"query_embedding": embedding_str, "top_k": top_k, "subject_id": subject_id}
        )
    else:
        sql_query = text("""
            SELECT 
                e.text_chunk,
                n.filename,
                s.name as subject_name,
                (1 - (e.embedding <=> CAST(:query_embedding AS vector))) as similarity
            FROM embeddings e
            JOIN notes n ON e.note_id = n.id
            JOIN subjects s ON n.subject_id = s.id
            ORDER BY e.embedding <=> CAST(:query_embedding AS vector)
            LIMIT :top_k
        """)
        result = db.execute(
            sql_query,
            {"query_embedding": embedding_str, "top_k": top_k}
        )
    
    results = []
    for row in result:
        results.append({
            "text": row.text_chunk,
            "filename": row.filename,
            "subject": row.subject_name,
            "similarity": float(row.similarity)
        })
    
    return results


def build_qa_chain():
    """
    Build a custom QA prompt template for the RAG system.
    Optimized for step-by-step mathematical problem solving.
    """
    template = """You are an expert AI Tutor specializing in academic subjects, particularly mathematics and technical topics.
Your role is to help students understand concepts and solve problems using the provided course materials.

CONTEXT FROM STUDENT'S NOTES:
{context}

STUDENT'S QUESTION:
{question}

RESPONSE GUIDELINES:

1. **ANSWER AVAILABILITY**:
   - Only use information from the provided context above
   - If the answer is not in the context, clearly state: "I cannot find this information in your uploaded notes."
   - Never make up information or use external knowledge not present in the context

2. **FOR MATHEMATICAL PROBLEMS** (Integration, Derivation, Equations, Calculations):
   - **Always provide step-by-step solutions**
   - Break down each step clearly with explanations
   - Use this format:
     
     **Step 1:** [Brief description of what you're doing]
     [Show the mathematical operation]
     
     **Step 2:** [Next step description]
     [Show the calculation]
     
     **Final Answer:** [Clearly state the final result]
   
   - Use LaTeX for ALL mathematical expressions:
     * Inline math: $expression$
     * Examples: $\\int x^2 dx$, $\\frac{{dy}}{{dx}}$, $e^x$, $\\sin(x)$, $\\sum_{{i=1}}^n$
   
   - For integration problems: Show integration steps, apply rules, simplify
   - For derivation problems: Show differentiation rules applied, chain rule if needed
   - For equation solving: Show algebraic manipulations clearly
   - Always verify and show your final answer

3. **FOR THEORY/CONCEPTUAL QUESTIONS**:
   - Provide clear, structured explanations
   - Use bullet points or numbered lists for clarity
   - Include relevant examples from the context
   - Define key terms when first mentioned

4. **FORMATTING**:
   - Use **bold** for important terms and step headers
   - Use bullet points or numbered lists for better readability
   - Keep paragraphs concise (2-3 sentences max)
   - Add line breaks between major sections

5. **TONE & STYLE**:
   - Professional and encouraging
   - Academic but easy to understand
   - Patient and thorough
   - Avoid being condescending

6. **EXAMPLES & CONTEXT**:
   - Reference specific theorems, formulas, or concepts from the context
   - If there are examples in the context, relate them to the question
   - Point out which part of the notes contains the relevant information

Now, provide a comprehensive answer to the student's question based on the context above:"""

    prompt = PromptTemplate(
        template=template,
        input_variables=["context", "question"]
    )
    
    return prompt


def ask_question(question: str, db: Session, subject_id: str = None) -> dict:
    """
    Answer a question using RAG pipeline.
    Optionally scoped to a specific subject.
    Returns answer and the sources used.
    """
    # Perform similarity search (scoped to subject if provided)
    search_results = similarity_search(question, db, top_k=3, subject_id=subject_id)
    
    # Check if we have any results
    if not search_results or len(search_results) == 0:
        return {
            "answer": "Answer not available in notes.",
            "sources": []
        }
    
    # Build context from search results
    context = "\n\n".join([
        f"From {result['filename']} ({result['subject']}):\n{result['text']}"
        for result in search_results
    ])
    
    # Build prompt
    prompt_template = build_qa_chain()
    prompt_text = prompt_template.format(context=context, question=question)
    
    # Generate answer using Gemini
    try:
        response = llm.invoke(prompt_text)
        answer = response.content
    except Exception as e:
        # If LLM fails, return the error
        answer = f"Error generating answer: {str(e)}"
    
    # Extract unique sources
    sources = list(set([
        f"{result['filename']} ({result['subject']})"
        for result in search_results
    ]))
    
    return {
        "answer": answer,
        "sources": sources
    }
